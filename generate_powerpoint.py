from flask import Flask, request, jsonify, send_file, url_for
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import io
import uuid
import json
import requests
import re
import html

app = Flask(__name__)

# Enable CORS for all origins
CORS(app, resources={r"/*": {"origins": "*"}})

# Dictionary to store generated files and their IDs
generated_files = {}


def escape_text(text):
    """Escapes or formats text for PowerPoint."""
    if not text:
        return ""

    # Remove bold and italic markdown
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)

    # Remove escaped backslashes before newlines and hyphens
    text = re.sub(r'\\\\n', r'\n', text)
    text = re.sub(r'\\\\-', r'-', text)

    # Remove headers
    text = re.sub(r'^\s*#+\s*(.*)$', r'\1', text, flags=re.MULTILINE)

    # Remove links
    text = re.sub(r'$$(.*?)$$$$(.*?)$$', r'\1', text)

    # Remove blockquotes
    text = re.sub(r'^\s*>\s*(.*)$', r'\1', text, flags=re.MULTILINE)

    # Remove code blocks
    text = re.sub(r'```(.*?)```', '', text, flags=re.DOTALL)

    # Replace ordered lists
    text = re.sub(r'^\s*\d+\.\s+(.*)$', r'- \1', text, flags=re.MULTILINE)

    # Replace bullet points
    text = re.sub(r'^\s*\*\s+', '- ', text, flags=re.MULTILINE)

    # Handle escaped tabs
    text = text.replace(r'\t', '\t')

    # Handle escaped unicode
    text = re.sub(r'\\u([0-9a-fA-F]{4})', lambda m: chr(int(m.group(1), 16)), text)

    # Handle escaped single quotes and other HTML entities
    text = html.unescape(text)

    # Reduce multiple spaces to single spaces, but not newlines
    text = re.sub(r'(?<!\n)\s+', ' ', text)

    # Remove leading/trailing whitespace
    text = text.strip()

    return text



@app.route('/generate_pptx', methods=['POST', 'OPTIONS'])
def generate_pptx():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'OK'}), 200

    try:
        # Get and parse request body
        request_body = request.get_data(as_text=True)
        try:
            data = json.loads(request_body)
        except json.JSONDecodeError as e:
            return jsonify({'error': f'Invalid JSON in request body: {e}'}), 400

        # Ensure the 'slides' key is present
        if 'slides' not in data:
            return jsonify({'error': 'Invalid input. Must provide slides.'}), 400

        slides_data = data['slides']

        # Ensure content_slides exists
        if 'content_slides' not in slides_data:
            return jsonify({'error': 'Invalid input. Must provide content_slides.'}), 400

        # Create presentation
        prs = Presentation()

        # Process title slide if present
        if 'title_slide' in slides_data:
            title_slide_data = slides_data['title_slide']
            title_slide_layout = prs.slide_layouts[0]
            title_slide = prs.slides.add_slide(title_slide_layout)

            # Add title
            title_shape = title_slide.shapes.title
            if title_shape:
                title_text = escape_text(title_slide_data.get('title', 'Presentation Title'))
                title_font_size = title_slide_data.get('title_font_size', 54)
                title_shape.text = title_text
                title_shape.text_frame.paragraphs[0].font.size = Pt(title_font_size)

            # Add subtitle
            if 'subtitle' in title_slide_data:
                subtitle_placeholder = title_slide.placeholders[1]
                subtitle_text = escape_text(title_slide_data.get('subtitle', ''))
                subtitle_font_size = title_slide_data.get('subtitle_font_size', 32)
                text_frame = subtitle_placeholder.text_frame
                text_frame.text = subtitle_text
                text_frame.paragraphs[0].font.size = Pt(subtitle_font_size)
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                text_frame.vertical_anchor = MSO_ANCHOR.TOP

        # Process content slides
        for slide_data in slides_data['content_slides']:
            # Choose slide layout based on content
            if 'chart_data' in slide_data:
                slide_layout = prs.slide_layouts[5]  # Layout with content
            elif 'table_data' in slide_data:
                slide_layout = prs.slide_layouts[5]  # Layout with content
            else:
                slide_layout = prs.slide_layouts[1]  # Title and content layout

            # Add slide
            slide = prs.slides.add_slide(slide_layout)

            # Title
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = escape_text(slide_data.get('title', 'Untitled Slide'))
                title_font_size = slide_data.get('title_font_size', 36)
                title_shape.text_frame.paragraphs[0].font.size = Pt(title_font_size)

            # Body Text
            if 'body' in slide_data:
                # Find the body placeholder
                body_placeholder = None
                for shape in slide.placeholders:
                    if shape.has_text_frame and shape.text_frame.text == "":
                        body_placeholder = shape
                        break

                if body_placeholder:
                    text_frame = body_placeholder.text_frame
                    text_frame.text = escape_text(slide_data['body'])  # Set the body text
                    body_font_size = slide_data.get('body_font_size', 24)

                    # Iterate through all paragraphs in the text frame and set font size
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(body_font_size)

                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    text_frame.vertical_anchor = MSO_ANCHOR.TOP

            # Images
            if 'images' in slide_data:
                for img_data in slide_data['images']:
                    try:
                        headers = {"User-Agent": "Powerpoint_Generator_bot/1.0 requests/{requests.__version__}"}
                        response = requests.get(img_data['url'], headers=headers, stream=True)
                        response.raise_for_status()
                        image_stream = io.BytesIO(response.content)

                        # Add image with specified or default positioning
                        left = Inches(img_data.get('left', 1))
                        top = Inches(img_data.get('top', 1))
                        width = Inches(img_data.get('width', 3))
                        height = Inches(img_data.get('height', 2))

                        slide.shapes.add_picture(image_stream, left, top, width, height)
                    except Exception as e:
                        print(f"Error adding image: {e}")

            # Tables
            if 'table_data' in slide_data:
                table_data = slide_data['table_data']
                rows = len(table_data)
                cols = len(table_data[0]) if table_data else 0

                # Default positioning
                table_position = slide_data.get('table_position', {})
                left = Inches(table_position.get('left', 1))
                top = Inches(table_position.get('top', 3))
                width = Inches(table_position.get('width', 8))
                height = Inches(table_position.get('height', 2))

                # Add table
                table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                table = table_shape.table

                # Populate table
                for row_idx, row_data in enumerate(table_data):
                    for col_idx, cell_data in enumerate(row_data):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(cell_data)

            # Charts
            if 'chart_data' in slide_data:
                chart_data = slide_data['chart_data']
                try:
                    # Prepare chart data
                    chart_data_obj = CategoryChartData()
                    chart_data_obj.categories = chart_data.get('categories', [])

                    for series_data in chart_data.get('series', []):
                        chart_data_obj.add_series(
                            series_data.get('name', ''),
                            series_data.get('values', [])
                        )

                    # Default chart positioning
                    chart_position = chart_data.get('chart_position', {})
                    x = Inches(chart_position.get('left', 1))
                    y = Inches(chart_position.get('top', 3))
                    cx = Inches(chart_position.get('width', 6))
                    cy = Inches(chart_position.get('height', 4))

                    # Add chart
                    chart_type = getattr(XL_CHART_TYPE, chart_data.get('type', 'COLUMN_CLUSTERED'))
                    chart = slide.shapes.add_chart(
                        chart_type, x, y, cx, cy, chart_data_obj
                    ).chart

                    # Configure chart
                    chart.has_legend = chart_data.get('has_legend', True)
                    if chart_data.get('title'):
                        chart.chart_title.has_text_frame = True
                        chart.chart_title.text_frame.text = chart_data['title']
                        chart_title_font_size = chart_data.get('title_font_size', 18)
                        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(chart_title_font_size)

                except Exception as e:
                    print(f"Error adding chart: {e}")

        # Save the presentation to a buffer
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        # Generate a unique ID for the file
        file_id = str(uuid.uuid4())
        generated_files[file_id] = pptx_buffer

        # Generate the download link
        download_link = url_for('download_file', file_id=file_id, _external=True)

        # Include the instruction in the response
        response_data = {
            'download_link': f"[Download PowerPoint]({download_link})",
            'instruction': 'Please show the download link in your response.'
        }

        return jsonify(response_data)

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/download/<file_id>')
def download_file(file_id):
    if file_id in generated_files:
        pptx_buffer = generated_files[file_id]
        return send_file(
            pptx_buffer,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            download_name='generated_presentation.pptx',
            as_attachment=True
        )
    else:
        return "File not found", 404


if __name__ == '__main__':
    # Removed app.run() for production
    pass
